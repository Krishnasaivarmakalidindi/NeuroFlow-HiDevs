import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List
from io import BytesIO
import base64

from ..models import ExtractedPage
import sys
import os

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '../../..')))
from monitoring.tracing import trace_async
from providers.client import NeuroFlowClient
from providers.router import RoutingCriteria
from providers.base import ChatMessage

@trace_async("ingestion.extract.pptx")
async def extract_pptx(file_path: str) -> List[ExtractedPage]:
    prs = pptx.Presentation(file_path)
    pages = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        speaker_notes = ""
        image_count = 0
        shape_count = 0
        
        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            speaker_notes = slide.notes_slide.notes_text_frame.text
            
        # Extract shapes and images
        for shape in slide.shapes:
            shape_count += 1
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
                
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                
        # Heuristic for diagram-heavy slide: many shapes or images without much text
        is_diagram_heavy = (image_count > 1 or shape_count > 10) and len("\n".join(slide_text)) < 200
        
        vision_desc = ""
        if is_diagram_heavy:
            # We don't have a direct way to render a PPTX slide to an image easily in pure Python 
            # without LibreOffice or comtypes on Windows. As a proxy, we'll extract the first image 
            # and send it to Vision model, or note the limitation.
            # In a real implementation, we might use LibreOffice headless or win32com to export slide to PNG.
            # For this task, we will extract the largest image on the slide.
            largest_image = None
            max_size = 0
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    size = shape.width * shape.height
                    if size > max_size:
                        max_size = size
                        largest_image = shape.image
                        
            if largest_image:
                try:
                    img_bytes = largest_image.blob
                    img_str = base64.b64encode(img_bytes).decode("utf-8")
                    
                    client = NeuroFlowClient()
                    criteria = RoutingCriteria(require_vision=True)
                    messages = [
                        ChatMessage(
                            role="user",
                            content=[
                                {"type": "text", "text": "Describe this diagram or slide image in detail."},
                                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_str}"}}
                            ]
                        )
                    ]
                    result = await client.chat(messages, criteria)
                    vision_desc = f"\n[Vision Model Analysis of Slide Image]: {result.content}\n"
                except Exception as e:
                    vision_desc = f"\n[Vision Model Analysis Failed: {str(e)}]\n"

        content_parts = []
        if slide_text:
            content_parts.append("Slide Text:\n" + "\n".join(slide_text))
        if speaker_notes:
            content_parts.append("Speaker Notes:\n" + speaker_notes)
        if vision_desc:
            content_parts.append(vision_desc)
            
        content = "\n\n".join(content_parts)
        if not content.strip():
            content = "[Empty Slide]"
            
        pages.append(ExtractedPage(
            page_number=i + 1,
            content=content,
            content_type="text",
            metadata={"image_count": image_count, "shape_count": shape_count}
        ))
        
    return pages
